from typing import List, Optional, Dict, TypedDict, Annotated, Sequence, TypeVar, Union
import PyPDF2
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.docstore.document import Document
from langchain_community.chat_models import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from duckduckgo_search import DDGS
import os
from business_term_sheet import BusinessTermSheet
from dotenv import load_dotenv
import json
from dataclasses import dataclass
import time
from requests.exceptions import RequestException

# Load environment variables
load_dotenv()

class FieldResult:
    def __init__(self, output_value: str, is_mentioned: bool, confidence: int, source: str):
        self.output_value = output_value
        self.is_mentioned = is_mentioned
        self.confidence = confidence
        self.source = source

    def to_dict(self):
        return {
            "output_value": self.output_value,
            "is_mentioned": self.is_mentioned,
            "confidence": self.confidence,
            "source": self.source
        }

class GraphState(TypedDict):
    field_name: str
    description: str
    context: str
    result: Optional[FieldResult]
    next_step: str

class PitchDeckAnalyzer:
    def __init__(self):
        self.embeddings = OpenAIEmbeddings()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        self.vector_store = None
        self.llm = ChatOpenAI(model="gpt-4", temperature=0)
        self.search = None  # Initialize DDGS for each search to avoid rate limits
        self.last_search_time = 0
        self.min_search_interval = 2  # Minimum seconds between searches

    def _safe_web_search(self, query: str, max_retries: int = 3) -> List[Dict]:
        """Perform web search with rate limiting and retries."""
        current_time = time.time()
        time_since_last_search = current_time - self.last_search_time
        
        if time_since_last_search < self.min_search_interval:
            time.sleep(self.min_search_interval - time_since_last_search)
        
        for attempt in range(max_retries):
            try:
                # Create a new DDGS instance for each search
                self.search = DDGS()
                results = list(self.search.text(query, max_results=2))
                self.last_search_time = time.time()
                return results
            except Exception as e:
                if "Ratelimit" in str(e) and attempt < max_retries - 1:
                    wait_time = (attempt + 1) * 5  # Exponential backoff
                    print(f"Rate limited, waiting {wait_time} seconds...")
                    time.sleep(wait_time)
                    continue
                raise e
        return []

    def extract_text_from_pdf(self, file_path: str) -> List[Document]:
        documents = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        doc = Document(
                            page_content=text,
                            metadata={"source": file_path, "page": page_num + 1}
                        )
                        documents.append(doc)
                print(f"Extracted {len(documents)} pages from PDF")
        except Exception as e:
            print(f"Error extracting text from PDF: {e}")
        return documents

    def extract_text_from_pptx(self, file_path: str) -> List[Document]:
        documents = []
        try:
            presentation = Presentation(file_path)
            for slide_num, slide in enumerate(presentation.slides):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                if text_content:
                    doc = Document(
                        page_content=" ".join(text_content),
                        metadata={"source": file_path, "slide": slide_num + 1}
                    )
                    documents.append(doc)
            print(f"Extracted {len(documents)} slides from PPTX")
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
        return documents

    def create_vector_store(self, file_path: str) -> None:
        try:
            if file_path.lower().endswith('.pdf'):
                documents = self.extract_text_from_pdf(file_path)
            elif file_path.lower().endswith('.pptx'):
                documents = self.extract_text_from_pptx(file_path)
            else:
                raise ValueError("Unsupported file format. Please provide a PDF or PPTX file.")

            if not documents:
                raise ValueError("No text content extracted from the document")

            split_docs = self.text_splitter.split_documents(documents)
            print(f"Created {len(split_docs)} text chunks")
            
            self.vector_store = FAISS.from_documents(split_docs, self.embeddings)
            print("Vector store created successfully")
            
        except Exception as e:
            print(f"Error creating vector store: {e}")
            raise

    def _calculate_confidence(self, similarity_scores: List[float] = None, source: str = "llm") -> int:
        """Calculate confidence score based on source and similarity scores."""
        if similarity_scores:
            # For vector store, use average similarity score
            avg_similarity = sum(similarity_scores) / len(similarity_scores)
            base_confidence = int(min(100, max(0, (1 - avg_similarity) * 100)))
        else:
            # Base confidence for different sources
            base_confidence = {
                "vector_store": 70,
                "web_search": 50,
                "llm": 30
            }.get(source, 0)
        
        # Apply source-specific modifiers
        if source == "web_search":
            # Reduce confidence for web search relative to vector store
            base_confidence = int(base_confidence * 0.8)
        elif source == "llm":
            # Further reduce confidence for LLM-only responses
            base_confidence = int(base_confidence * 0.6)
        
        return min(100, max(0, base_confidence))

    def query_field(self, field_name: str, description: str) -> FieldResult:
        """Query a single field using sequential fallback strategy."""
        try:
            # Try vector store first
            if self.vector_store:
                query = f"Find information about {field_name}: {description}"
                docs_with_scores = self.vector_store.similarity_search_with_score(query, k=3)
                
                if docs_with_scores:
                    similarity_scores = [score for _, score in docs_with_scores]
                    confidence = self._calculate_confidence(similarity_scores, "vector_store")
                    
                    if confidence >= 30:
                        context = "\n".join([doc.page_content for doc, _ in docs_with_scores])
                        result = self._extract_with_llm(field_name, description, context)
                        if result.is_mentioned:
                            result.confidence = confidence
                            result.source = "vector_store"
                            return result

            # Try web search if vector store fails
            try:
                # Use a shorter query for web search
                search_query = f"{field_name} {description.split('.')[0]}"  # Only use the first sentence
                search_results = self._safe_web_search(search_query)
                
                if search_results:
                    context = "\n".join([f"{result['title']}: {result['body']}" for result in search_results])
                    result = self._extract_with_llm(field_name, description, context)
                    if result.is_mentioned:
                        result.confidence = self._calculate_confidence(source="web_search")
                        result.source = "web_search"
                        return result
            except Exception as e:
                print(f"Web search failed: {e}")

            # Fallback to LLM generation
            prompt = f"""Generate a response for the following field from a pitch deck:
            Field: {field_name}
            Description: {description}
            
            If you cannot provide accurate information, respond with 'Not mentioned'.
            """

            response = self.llm.invoke([
                SystemMessage(content="You are a precise business analyst. Generate information only if you are confident."),
                HumanMessage(content=prompt)
            ])
            
            return FieldResult(
                output_value=response.content.strip(),
                is_mentioned=response.content.strip() != "Not mentioned",
                confidence=self._calculate_confidence(source="llm"),
                source="llm_generation"
            )

        except Exception as e:
            print(f"Error querying field {field_name}: {e}")
            return FieldResult(
                output_value="Error during extraction",
                is_mentioned=False,
                confidence=self._calculate_confidence(),  # Default to 0
                source="error"
            )

    def _extract_with_llm(self, field_name: str, description: str, context: str) -> FieldResult:
        """Helper method to extract information using LLM from context."""
        prompt = f"""Based on the following context, extract the {field_name}. 
        The field description is: {description}
        If the information is not explicitly mentioned, respond with 'Not mentioned'.
        
        Context:
        {context}
        
        {field_name}:"""

        response = self.llm.invoke([
            SystemMessage(content="You are a precise business analyst. Extract specific information from the given context."),
            HumanMessage(content=prompt)
        ])
        
        extracted_value = response.content.strip()
        return FieldResult(
            output_value=extracted_value,
            is_mentioned=extracted_value != "Not mentioned",
            confidence=0,  # Will be set by the caller
            source=""  # Will be set by the caller
        )

    def analyze_pitch_deck(self, file_path: str) -> Dict[str, Dict[str, any]]:
        """Analyze the entire pitch deck."""
        try:
            print(f"\nAnalyzing pitch deck: {file_path}")
            self.create_vector_store(file_path)
            
            extracted_data = {}
            for field_name, field_info in BusinessTermSheet.model_fields.items():
                description = field_info.description or field_name
                print(f"\nExtracting {field_name}...")
                result = self.query_field(field_name, description)
                extracted_data[field_name] = result.to_dict()
            
            return extracted_data
            
        except Exception as e:
            print(f"Error analyzing pitch deck: {e}")
            raise

def main():
    analyzer = PitchDeckAnalyzer()
    file_path = "/Users/chandima/repos/resume-analysis/term-sheets/data/pitch/BISTEC Care Deck Short version.pdf"
    
    try:
        results = analyzer.analyze_pitch_deck(file_path)
        print("\nExtracted Business Terms:")
        print(json.dumps(results, indent=2))
    except Exception as e:
        print(f"Error analyzing pitch deck: {e}")

if __name__ == "__main__":
    main()
